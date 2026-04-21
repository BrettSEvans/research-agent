"""
Format conversion utilities: convert various file formats to PDF for extraction.

Supports: .pptx, .ppt, .docx, .doc, .xlsx, and Google Slides (via OAuth2).
All conversions output PDF files for the existing extraction pipeline.
"""

import os
import subprocess
import tempfile
from pathlib import Path
from io import BytesIO

from pptx import Presentation as PPTXPresentation
from docx import Document as DocxDocument
from openpyxl import load_workbook
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.units import inch
from PIL import Image
import pdf2image
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google.auth.transport.requests import Request
from tenacity import retry, stop_after_attempt, wait_fixed, retry_if_exception_type


def detect_file_type(filename: str) -> str:
    """Detect file type from extension.

    Args:
        filename: Name of the uploaded file

    Returns: File type string ('pdf', 'pptx', 'ppt', 'docx', 'doc', 'xlsx')

    Raises:
        ValueError: If file extension is not supported
    """
    if not filename:
        raise ValueError("Empty filename")

    ext = filename.lower().rsplit('.', 1)[-1]

    supported = {
        'pdf': 'pdf',
        'pptx': 'pptx',
        'ppt': 'ppt',
        'docx': 'docx',
        'doc': 'doc',
        'xlsx': 'xlsx',
    }

    if ext not in supported:
        raise ValueError(
            f"Unsupported file format: .{ext}. "
            f"Supported formats: {', '.join(supported.keys())}"
        )

    return ext


@retry(
    stop=stop_after_attempt(2),
    wait=wait_fixed(2),
    retry=retry_if_exception_type(subprocess.TimeoutExpired),
    reraise=True,
)
def _run_soffice(cmd: list, timeout: int = 60) -> subprocess.CompletedProcess:
    """Run LibreOffice command with timeout retry.

    Retries on subprocess.TimeoutExpired with 2s fixed delay (2 attempts total).
    Other exceptions are not retried.

    Args:
        cmd: Command list for subprocess.run
        timeout: Timeout in seconds (default 60)

    Returns:
        subprocess.CompletedProcess result

    Raises:
        subprocess.TimeoutExpired: If timeout occurs on final attempt
        Any other exception from subprocess.run
    """
    return subprocess.run(cmd, capture_output=True, timeout=timeout)


async def convert_to_pdf(file_path: Path, file_type: str) -> Path:
    """Convert a document to PDF. Returns path to the generated PDF.

    Args:
        file_path: Path to the uploaded file (e.g., /tmp/upload_xyz.pptx)
        file_type: File type from detect_file_type()

    Returns: Path to the generated PDF (same directory as input, .pdf extension)

    Raises:
        ValueError: If conversion fails
    """
    pdf_path = file_path.with_suffix('.pdf')

    try:
        if file_type == 'pdf':
            # PDF already, just move to output path
            file_path.rename(pdf_path)
            return pdf_path

        elif file_type == 'pptx':
            convert_pptx_to_pdf(file_path, pdf_path)

        elif file_type == 'ppt':
            convert_ppt_to_pdf(file_path, pdf_path)

        elif file_type == 'docx':
            convert_docx_to_pdf(file_path, pdf_path)

        elif file_type == 'doc':
            convert_doc_to_pdf(file_path, pdf_path)

        elif file_type == 'xlsx':
            convert_xlsx_to_pdf(file_path, pdf_path)

        else:
            raise ValueError(f"Unknown file type: {file_type}")

        # Clean up original file
        if file_path.exists():
            file_path.unlink()

        if not pdf_path.exists():
            raise ValueError("PDF conversion produced no output file")

        return pdf_path

    except Exception as e:
        # Clean up on error
        if file_path.exists():
            file_path.unlink()
        if pdf_path.exists():
            pdf_path.unlink()
        raise ValueError(f"Failed to convert {file_type.upper()}: {str(e)}")


def convert_pptx_to_pdf(pptx_path: Path, pdf_path: Path) -> None:
    """Convert PowerPoint (.pptx) to PDF by rendering slides as images."""
    try:
        # Load presentation
        prs = PPTXPresentation(str(pptx_path))

        # Convert each slide to image and stack into PDF
        images = []
        for slide in prs.slides:
            # Save slide as temporary PNG
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                # python-pptx doesn't render directly, so we'll use a workaround:
                # For now, create a simple text-based PDF with slide content
                # In production, consider using LibreOffice for higher quality
                pass

        # Fallback: use LibreOffice if available (with retry on timeout)
        result = _run_soffice(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir',
             str(pdf_path.parent), str(pptx_path)],
            timeout=60
        )

        if result.returncode != 0:
            # Fallback: create a simple text PDF from slide text
            _pptx_to_text_pdf(prs, pdf_path)

    except Exception as e:
        raise ValueError(f"PPTX conversion failed: {str(e)}")


def _pptx_to_text_pdf(prs, pdf_path: Path) -> None:
    """Create PDF from PPTX text content (when LibreOffice unavailable)."""
    try:
        c = canvas.Canvas(str(pdf_path), pagesize=letter)
        width, height = letter

        for slide_num, slide in enumerate(prs.slides, 1):
            # Add slide number
            c.setFont("Helvetica-Bold", 14)
            c.drawString(0.5 * inch, height - 0.5 * inch, f"Slide {slide_num}")

            # Extract and draw text from shapes
            y_pos = height - 1.0 * inch
            c.setFont("Helvetica", 11)

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    # Wrap text and draw
                    for line in shape.text.split('\n'):
                        if y_pos < 0.5 * inch:
                            c.showPage()
                            y_pos = height - 0.5 * inch
                        c.drawString(0.75 * inch, y_pos, line[:80])  # Truncate long lines
                        y_pos -= 0.25 * inch

            c.showPage()

        c.save()
    except Exception as e:
        raise ValueError(f"Failed to create text PDF from PPTX: {str(e)}")


def convert_ppt_to_pdf(ppt_path: Path, pdf_path: Path) -> None:
    """Convert legacy PowerPoint (.ppt) to PDF using LibreOffice.

    Retries on timeout with 2 attempts via Tenacity (2s fixed delay between attempts).
    """
    try:
        result = _run_soffice(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir',
             str(pdf_path.parent), str(ppt_path)],
            timeout=60
        )

        if result.returncode != 0:
            raise ValueError(f"LibreOffice conversion failed: {result.stderr.decode()}")

        # LibreOffice creates PDF with original filename
        libreoffice_pdf = pdf_path.parent / ppt_path.with_suffix('.pdf').name
        if libreoffice_pdf.exists() and libreoffice_pdf != pdf_path:
            libreoffice_pdf.rename(pdf_path)

    except FileNotFoundError:
        raise ValueError(
            "LibreOffice (soffice) not found. Required for .ppt conversion. "
            "Install with: apt-get install libreoffice"
        )
    except Exception as e:
        raise ValueError(f"PPT conversion failed: {str(e)}")


def convert_docx_to_pdf(docx_path: Path, pdf_path: Path) -> None:
    """Convert Word (.docx) to PDF by extracting text and formatting."""
    try:
        doc = DocxDocument(str(docx_path))

        # Create PDF with document content
        c = canvas.Canvas(str(pdf_path), pagesize=letter)
        width, height = letter
        y_pos = height - 0.5 * inch

        c.setFont("Helvetica", 11)

        for para in doc.paragraphs:
            if para.text:
                # Handle line wrapping for long text
                lines = _wrap_text(para.text, max_width=80)
                for line in lines:
                    if y_pos < 0.5 * inch:
                        c.showPage()
                        y_pos = height - 0.5 * inch
                    c.drawString(0.75 * inch, y_pos, line)
                    y_pos -= 0.25 * inch

        # Add tables
        for table in doc.tables:
            y_pos -= 0.25 * inch  # Space before table
            for row in table.rows:
                for cell in row.cells:
                    if y_pos < 0.5 * inch:
                        c.showPage()
                        y_pos = height - 0.5 * inch
                    c.drawString(0.75 * inch, y_pos, cell.text[:60])
                y_pos -= 0.2 * inch

        c.save()

    except Exception as e:
        raise ValueError(f"DOCX conversion failed: {str(e)}")


def convert_doc_to_pdf(doc_path: Path, pdf_path: Path) -> None:
    """Convert legacy Word (.doc) to PDF using LibreOffice.

    Retries on timeout with 2 attempts via Tenacity (2s fixed delay between attempts).
    """
    try:
        result = _run_soffice(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir',
             str(pdf_path.parent), str(doc_path)],
            timeout=60
        )

        if result.returncode != 0:
            raise ValueError(f"LibreOffice conversion failed: {result.stderr.decode()}")

        # LibreOffice creates PDF with original filename
        libreoffice_pdf = pdf_path.parent / doc_path.with_suffix('.pdf').name
        if libreoffice_pdf.exists() and libreoffice_pdf != pdf_path:
            libreoffice_pdf.rename(pdf_path)

    except FileNotFoundError:
        raise ValueError(
            "LibreOffice (soffice) not found. Required for .doc conversion. "
            "Install with: apt-get install libreoffice"
        )
    except Exception as e:
        raise ValueError(f"DOC conversion failed: {str(e)}")


def convert_xlsx_to_pdf(xlsx_path: Path, pdf_path: Path) -> None:
    """Convert Excel (.xlsx) to PDF by rendering sheets as text."""
    try:
        wb = load_workbook(str(xlsx_path))

        c = canvas.Canvas(str(pdf_path), pagesize=letter)
        width, height = letter
        y_pos = height - 0.5 * inch

        for sheet in wb.sheetnames:
            ws = wb[sheet]

            # Sheet title
            c.setFont("Helvetica-Bold", 12)
            if y_pos < 1.0 * inch:
                c.showPage()
                y_pos = height - 0.5 * inch
            c.drawString(0.5 * inch, y_pos, f"Sheet: {sheet}")
            y_pos -= 0.3 * inch

            # Sheet data
            c.setFont("Helvetica", 9)
            for row in ws.iter_rows(min_row=1, max_row=100, values_only=True):  # Limit rows
                if y_pos < 0.5 * inch:
                    c.showPage()
                    y_pos = height - 0.5 * inch

                row_text = " | ".join(str(cell) if cell else "" for cell in row)
                row_text = row_text[:90]  # Truncate long rows
                c.drawString(0.75 * inch, y_pos, row_text)
                y_pos -= 0.2 * inch

        c.save()

    except Exception as e:
        raise ValueError(f"XLSX conversion failed: {str(e)}")


def _wrap_text(text: str, max_width: int = 80) -> list:
    """Simple text wrapping."""
    if len(text) <= max_width:
        return [text]

    lines = []
    current = ""
    for word in text.split():
        if len(current) + len(word) + 1 <= max_width:
            current += (" " if current else "") + word
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


async def fetch_google_slides_pdf(access_token: str, presentation_id: str, pdf_path: Path) -> None:
    """Download PDF export of a Google Slides presentation using OAuth2 token.

    Args:
        access_token: Google OAuth2 access token with Google Drive access
        presentation_id: Google Slides presentation ID (from URL)
        pdf_path: Path where to save the downloaded PDF

    Raises:
        ValueError: If download fails
    """
    try:
        from google.oauth2.credentials import Credentials

        # Create credentials from access token
        credentials = Credentials(token=access_token)

        # Build Google Drive API client
        drive_service = build('drive', 'v3', credentials=credentials)

        # Export presentation as PDF
        request = drive_service.files().export_media(
            fileId=presentation_id,
            mimeType='application/pdf'
        )

        # Download PDF to file
        with open(pdf_path, 'wb') as fh:
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                status, done = downloader.next_chunk()

        if not pdf_path.exists() or pdf_path.stat().st_size == 0:
            raise ValueError("Google Slides PDF export resulted in empty file")

    except Exception as e:
        if pdf_path.exists():
            pdf_path.unlink()
        raise ValueError(f"Failed to download Google Slides PDF: {str(e)}")
